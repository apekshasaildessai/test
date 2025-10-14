#!/usr/bin/env python3
"""
Credora Capital PowerPoint Presentation Generator
Creates a native PPTX file with professional branding and all investor information.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Credora Capital Brand Colors
GOLD_COLOR = RGBColor(212, 175, 55)  # #D4AF37
DARK_GRAY = RGBColor(44, 44, 44)     # #2c2c2c
WHITE_COLOR = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)

def create_credora_presentation():
    """Create the complete Credora Capital presentation"""
    
    # Create presentation object
    prs = Presentation()
    
    # Set slide size to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    create_title_slide(prs)
    
    # Slide 2: Company Overview
    create_company_overview_slide(prs)
    
    # Slide 3: Fund Details
    create_fund_details_slide(prs)
    
    # Slide 4: Investment Strategy
    create_investment_strategy_slide(prs)
    
    # Slide 5: Investment Objective
    create_investment_objective_slide(prs)
    
    # Slide 6: Market Context
    create_market_context_slide(prs)
    
    # Slide 7: Ideal Investor Profile
    create_investor_profile_slide(prs)
    
    # Slide 8: Key Metrics
    create_key_metrics_slide(prs)
    
    # Slide 9: Risk Factors
    create_risk_factors_slide(prs)
    
    # Slide 10: J-Curve & Capital Deployment
    create_j_curve_slide(prs)
    
    # Slide 11: Fund Management Team
    create_management_team_slide(prs)
    
    # Slide 12: Future Team Expansion
    create_team_expansion_slide(prs)
    
    # Slide 13: Fees & Terms
    create_fees_terms_slide(prs)
    
    # Slide 14: Investor Eligibility
    create_investor_eligibility_slide(prs)
    
    # Slide 15: Thank You
    create_thank_you_slide(prs)
    
    # Save presentation
    output_file = '/app/Credora_Capital_Investor_Presentation.pptx'
    prs.save(output_file)
    print(f"✅ PowerPoint presentation saved as: {output_file}")
    return output_file

def add_logo_to_slide(slide, x=Inches(11.5), y=Inches(0.3), width=Inches(1.2)):
    """Add Credora Capital logo to slide"""
    try:
        logo_path = '/app/credora_logo.jpeg'
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, x, y, width=width)
    except Exception as e:
        print(f"Could not add logo: {e}")

def set_slide_background(slide, color=DARK_GRAY):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_title_slide(prs):
    """Create title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    set_slide_background(slide)
    
    # Add main logo (larger, centered)
    try:
        logo_path = '/app/credora_logo.jpeg'
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(5.5), Inches(1.5), width=Inches(2.5))
    except Exception as e:
        print(f"Could not add main logo: {e}")
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Investor Information Document"
    
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_font = title_para.font
    title_font.name = 'Segoe UI'
    title_font.size = Pt(48)
    title_font.color.rgb = GOLD_COLOR
    title_font.bold = True
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Unlocking India's Emerging Equity Potential"
    
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    
    # Year
    year_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.33), Inches(0.5))
    year_frame = year_box.text_frame
    year_frame.text = "2024"
    
    year_para = year_frame.paragraphs[0]
    year_para.alignment = PP_ALIGN.CENTER
    year_font = year_para.font
    year_font.name = 'Segoe UI'
    year_font.size = Pt(18)
    year_font.color.rgb = GOLD_COLOR

def create_standard_slide(prs, title_text):
    """Create a standard slide with header"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    set_slide_background(slide)
    
    # Add logo
    add_logo_to_slide(slide)
    
    # Header with golden line
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10.5), Inches(1))
    header_frame = header_box.text_frame
    header_frame.text = title_text
    
    header_para = header_frame.paragraphs[0]
    header_font = header_para.font
    header_font.name = 'Segoe UI'
    header_font.size = Pt(32)
    header_font.color.rgb = GOLD_COLOR
    header_font.bold = True
    
    # Add golden line under header
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.2), Inches(10.5), Inches(1.2))
    line.line.color.rgb = GOLD_COLOR
    line.line.width = Pt(3)
    
    return slide

def add_bullet_points(slide, points, start_y=Inches(2.5), start_x=Inches(1)):
    """Add bullet points to slide"""
    for i, point in enumerate(points):
        y_pos = start_y + Inches(0.6 * i)
        
        # Bullet symbol
        bullet_box = slide.shapes.add_textbox(start_x, y_pos, Inches(0.3), Inches(0.4))
        bullet_frame = bullet_box.text_frame
        bullet_frame.text = "•"
        bullet_para = bullet_frame.paragraphs[0]
        bullet_font = bullet_para.font
        bullet_font.name = 'Segoe UI'
        bullet_font.size = Pt(16)
        bullet_font.color.rgb = GOLD_COLOR
        bullet_font.bold = True
        
        # Point text
        text_box = slide.shapes.add_textbox(start_x + Inches(0.4), y_pos, Inches(10), Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = point
        text_para = text_frame.paragraphs[0]
        text_font = text_para.font
        text_font.name = 'Segoe UI'
        text_font.size = Pt(16)
        text_font.color.rgb = WHITE_COLOR

def create_company_overview_slide(prs):
    """Create company overview slide"""
    slide = create_standard_slide(prs, "Company Overview")
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Greenflow Ventures Pvt Ltd"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    subtitle_font.bold = True
    
    # Bullet points
    points = [
        "Positioned as a SEBI-compliant alternative investment platform",
        "Focus: Unlocking India's emerging equity potential through structured, data-driven, and regulation-aligned investments",
        "Flagship Fund: Credora Capital (Category I AIF)",
        "Outlook: Promising, aiming to tap into India's vibrant primary market ecosystem"
    ]
    
    add_bullet_points(slide, points)

def create_fund_details_slide(prs):
    """Create fund details slide"""
    slide = create_standard_slide(prs, "Fund Details")
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Credora Capital"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    subtitle_font.bold = True
    
    # Create info boxes
    info_items = [
        ("Fund Name", "Credora Capital"),
        ("Type", "Category I AIF (under SEBI AIF Regulations, 2012)"),
        ("Structure", "Close-ended fund"),
        ("Registration", "To be registered with SEBI as per AIF norms")
    ]
    
    for i, (label, value) in enumerate(info_items):
        x_pos = Inches(1) if i % 2 == 0 else Inches(6.5)
        y_pos = Inches(2.8) if i < 2 else Inches(4.5)
        
        # Create box background
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_pos, Inches(5), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(60, 60, 60)
        box.line.color.rgb = GOLD_COLOR
        box.line.width = Pt(2)
        
        # Label text
        label_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(0.1), Inches(4.6), Inches(0.4))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_font = label_para.font
        label_font.name = 'Segoe UI'
        label_font.size = Pt(14)
        label_font.color.rgb = GOLD_COLOR
        label_font.bold = True
        
        # Value text
        value_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(0.5), Inches(4.6), Inches(0.6))
        value_frame = value_box.text_frame
        value_frame.text = value
        value_para = value_frame.paragraphs[0]
        value_font = value_para.font
        value_font.name = 'Segoe UI'
        value_font.size = Pt(12)
        value_font.color.rgb = WHITE_COLOR

def create_investment_strategy_slide(prs):
    """Create investment strategy slide"""
    slide = create_standard_slide(prs, "Investment Strategy")
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Capitalizing on India's Primary Market Growth"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    subtitle_font.bold = True
    
    # Two columns
    # Left column - Key Avenues
    left_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(5.5), Inches(3))
    left_frame = left_box.text_frame
    left_frame.text = "Key Avenues\n• Anchor investing in IPOs\n• Pre-IPO and early-stage investment rounds"
    
    for i, para in enumerate(left_frame.paragraphs):
        if i == 0:  # Header
            para.font.name = 'Segoe UI'
            para.font.size = Pt(18)
            para.font.color.rgb = GOLD_COLOR
            para.font.bold = True
        else:  # Bullet points
            para.font.name = 'Segoe UI'
            para.font.size = Pt(14)
            para.font.color.rgb = WHITE_COLOR
    
    # Right column - Fund Aims
    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.5), Inches(5.5), Inches(3))
    right_frame = right_box.text_frame
    right_frame.text = "Fund Aims\n• Maximize alpha by investing before public price discovery\n• Optimize liquidity through participation in multiple IPO rounds\n• Support SMEs with growth capital in their pre-IPO phases"
    
    for i, para in enumerate(right_frame.paragraphs):
        if i == 0:  # Header
            para.font.name = 'Segoe UI'
            para.font.size = Pt(18)
            para.font.color.rgb = GOLD_COLOR
            para.font.bold = True
        else:  # Bullet points
            para.font.name = 'Segoe UI'
            para.font.size = Pt(14)
            para.font.color.rgb = WHITE_COLOR
    
    # Highlight box
    highlight_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.8), Inches(11.33), Inches(1))
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = GOLD_COLOR
    highlight_box.line.width = Pt(0)
    
    highlight_text = slide.shapes.add_textbox(Inches(1.2), Inches(6), Inches(11), Inches(0.6))
    highlight_frame = highlight_text.text_frame
    highlight_frame.text = "Hybrid Strategy: Combines risk-mitigated anchor allocations with strategic exposure to high-growth early-stage companies"
    highlight_para = highlight_frame.paragraphs[0]
    highlight_para.alignment = PP_ALIGN.CENTER
    highlight_font = highlight_para.font
    highlight_font.name = 'Segoe UI'
    highlight_font.size = Pt(14)
    highlight_font.color.rgb = DARK_GRAY
    highlight_font.bold = True

def create_investment_objective_slide(prs):
    """Create investment objective slide"""
    slide = create_standard_slide(prs, "Investment Objective")
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Generating Consistent & Risk-Adjusted Returns"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    subtitle_font.bold = True
    
    # Four objective boxes
    objectives = [
        ("Anchor Investments in Mainboard IPOs", "Leverage preferential allocations for early exposure, listing gains, and post-lock-in upside"),
        ("SME IPO Participation", "Focus on SMEs with high-growth potential and less institutional coverage"),
        ("Pre-IPO Rounds", "Back scalable businesses 12-36 months from listing"),
        ("Early-Stage Startups", "Opportunistic allocation to startups with strong teams and near-term listing roadmaps")
    ]
    
    for i, (title, desc) in enumerate(objectives):
        x_pos = Inches(1) if i % 2 == 0 else Inches(7)
        y_pos = Inches(2.8) if i < 2 else Inches(4.8)
        
        # Create box
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_pos, Inches(5.5), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(60, 60, 60)
        box.line.color.rgb = GOLD_COLOR
        box.line.width = Pt(2)
        
        # Title
        title_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(0.1), Inches(5.1), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_font = title_para.font
        title_font.name = 'Segoe UI'
        title_font.size = Pt(14)
        title_font.color.rgb = GOLD_COLOR
        title_font.bold = True
        
        # Description
        desc_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(0.6), Inches(5.1), Inches(0.8))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_para = desc_frame.paragraphs[0]
        desc_font = desc_para.font
        desc_font.name = 'Segoe UI'
        desc_font.size = Pt(11)
        desc_font.color.rgb = WHITE_COLOR

def create_market_context_slide(prs):
    """Create market context slide"""
    slide = create_standard_slide(prs, "Market Context & Rationale")
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "India's Capital Markets Transformation"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    subtitle_font.bold = True
    
    # Stats boxes
    stats = [
        ("400+", "IPOs launched since FY 2020"),
        ("₹1.25 Lakh Cr", "Capital raised"),
        ("25%+", "CAGR in SME listings (5 years)")
    ]
    
    for i, (number, label) in enumerate(stats):
        x_pos = Inches(1 + i * 4)
        
        # Stat box
        stat_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, Inches(2.8), Inches(3.5), Inches(1.2))
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = GOLD_COLOR
        stat_box.line.width = Pt(0)
        
        # Number
        num_box = slide.shapes.add_textbox(x_pos, Inches(2.9), Inches(3.5), Inches(0.6))
        num_frame = num_box.text_frame
        num_frame.text = number
        num_para = num_frame.paragraphs[0]
        num_para.alignment = PP_ALIGN.CENTER
        num_font = num_para.font
        num_font.name = 'Segoe UI'
        num_font.size = Pt(24)
        num_font.color.rgb = DARK_GRAY
        num_font.bold = True
        
        # Label
        label_box = slide.shapes.add_textbox(x_pos, Inches(3.5), Inches(3.5), Inches(0.4))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.alignment = PP_ALIGN.CENTER
        label_font = label_para.font
        label_font.name = 'Segoe UI'
        label_font.size = Pt(12)
        label_font.color.rgb = DARK_GRAY
        label_font.bold = True
    
    # Key points
    points = [
        "Anchor Investor Edge: Consistently outperformed due to better entry price and insider participation",
        "Retail Interest Growing: Institutions with early entry rights maintain a substantial edge",
        "Credora Capital's Role: Bridging the capital gap for SMEs and leveraging this trend through professional primary market investment management"
    ]
    
    add_bullet_points(slide, points, Inches(4.5))

# Continue with remaining slides...
def create_investor_profile_slide(prs):
    """Create ideal investor profile slide"""
    slide = create_standard_slide(prs, "Ideal Investor Profile")
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Who is this Fund For?"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    subtitle_font.bold = True
    
    points = [
        "Investors seeking moderate to high-risk investments in India's high-growth equity",
        "Those with a medium-to-long term investment horizon (5+ years)",
        "Investors who understand and accept market volatility and potential lock-in",
        "Seeking diversified, actively managed exposure to primary market opportunities"
    ]
    
    add_bullet_points(slide, points)
    
    # Investor types
    types_title = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(10), Inches(0.4))
    types_title_frame = types_title.text_frame
    types_title_frame.text = "Investor Types:"
    types_title_para = types_title_frame.paragraphs[0]
    types_title_font = types_title_para.font
    types_title_font.name = 'Segoe UI'
    types_title_font.size = Pt(16)
    types_title_font.color.rgb = GOLD_COLOR
    types_title_font.bold = True
    
    types_text = "QIBs • HNIs • Family Offices • Corporates • Investment Trusts"
    types_box = slide.shapes.add_textbox(Inches(1), Inches(5.6), Inches(11), Inches(0.4))
    types_box_frame = types_box.text_frame
    types_box_frame.text = types_text
    types_box_para = types_box_frame.paragraphs[0]
    types_box_font = types_box_para.font
    types_box_font.name = 'Segoe UI'
    types_box_font.size = Pt(14)
    types_box_font.color.rgb = WHITE_COLOR

def create_key_metrics_slide(prs):
    """Create key metrics slide"""
    slide = create_standard_slide(prs, "Key Metrics & Fund Structure")
    
    # Highlight metric
    highlight_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(2), Inches(7.33), Inches(1.5))
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = GOLD_COLOR
    highlight_box.line.width = Pt(0)
    
    # Main value
    value_box = slide.shapes.add_textbox(Inches(3), Inches(2.1), Inches(7.33), Inches(0.7))
    value_frame = value_box.text_frame
    value_frame.text = "₹1000 Crores"
    value_para = value_frame.paragraphs[0]
    value_para.alignment = PP_ALIGN.CENTER
    value_font = value_para.font
    value_font.name = 'Segoe UI'
    value_font.size = Pt(36)
    value_font.color.rgb = DARK_GRAY
    value_font.bold = True
    
    # Label
    label_box = slide.shapes.add_textbox(Inches(3), Inches(2.8), Inches(7.33), Inches(0.4))
    label_frame = label_box.text_frame
    label_frame.text = "Minimum AUM Target (Initial Close)"
    label_para = label_frame.paragraphs[0]
    label_para.alignment = PP_ALIGN.CENTER
    label_font = label_para.font
    label_font.name = 'Segoe UI'
    label_font.size = Pt(16)
    label_font.color.rgb = DARK_GRAY
    label_font.bold = True
    
    # Additional metrics
    metrics = [
        ("Fund Tenure", "7 to 8 years (with possibility of extension by LP consent)"),
        ("Minimum Investment", "₹1 crore (aligned with SEBI's Qualified Investor norms)"),
        ("Fund Type", "Category I AIF (SEBI AIF Regulations, 2012)")
    ]
    
    for i, (metric_label, metric_value) in enumerate(metrics):
        y_pos = Inches(4.2 + i * 0.8)
        
        # Label
        metric_label_box = slide.shapes.add_textbox(Inches(1), y_pos, Inches(3), Inches(0.4))
        metric_label_frame = metric_label_box.text_frame
        metric_label_frame.text = metric_label + ":"
        metric_label_para = metric_label_frame.paragraphs[0]
        metric_label_font = metric_label_para.font
        metric_label_font.name = 'Segoe UI'
        metric_label_font.size = Pt(14)
        metric_label_font.color.rgb = GOLD_COLOR
        metric_label_font.bold = True
        
        # Value
        metric_value_box = slide.shapes.add_textbox(Inches(4.5), y_pos, Inches(7), Inches(0.4))
        metric_value_frame = metric_value_box.text_frame
        metric_value_frame.text = metric_value
        metric_value_para = metric_value_frame.paragraphs[0]
        metric_value_font = metric_value_para.font
        metric_value_font.name = 'Segoe UI'
        metric_value_font.size = Pt(14)
        metric_value_font.color.rgb = WHITE_COLOR

def create_risk_factors_slide(prs):
    """Create risk factors slide"""
    slide = create_standard_slide(prs, "Risk Factors")
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Understanding the Investment Landscape"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    subtitle_font.bold = True
    
    # Risk categories
    risks = [
        ("Market Risk", "Equity investments are subject to market volatility"),
        ("Liquidity Risk", "Pre-IPO and SME investments may have limited liquidity"),
        ("Regulatory Risk", "Changes in SEBI guidelines, taxation, or government regulations"),
        ("Business Risk", "Higher operational/financial risks in SMEs and early-stage companies"),
        ("Concentration Risk", "Potential for concentrated exposure in select sectors/companies"),
        ("Valuation Risk", "Uncertainties in pre-IPO valuations due to limited comparables")
    ]
    
    for i, (risk_type, risk_desc) in enumerate(risks):
        x_pos = Inches(1) if i % 2 == 0 else Inches(7)
        y_pos = Inches(2.8 + (i // 2) * 1.2)
        
        # Risk box
        risk_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_pos, Inches(5.5), Inches(1))
        risk_box.fill.solid()
        risk_box.fill.fore_color.rgb = RGBColor(80, 40, 40)
        risk_box.line.color.rgb = RGBColor(200, 80, 80)
        risk_box.line.width = Pt(2)
        
        # Risk type
        type_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(0.1), Inches(5.1), Inches(0.3))
        type_frame = type_box.text_frame
        type_frame.text = risk_type
        type_para = type_frame.paragraphs[0]
        type_font = type_para.font
        type_font.name = 'Segoe UI'
        type_font.size = Pt(14)
        type_font.color.rgb = RGBColor(255, 120, 120)
        type_font.bold = True
        
        # Risk description
        desc_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(0.4), Inches(5.1), Inches(0.5))
        desc_frame = desc_box.text_frame
        desc_frame.text = risk_desc
        desc_para = desc_frame.paragraphs[0]
        desc_font = desc_para.font
        desc_font.name = 'Segoe UI'
        desc_font.size = Pt(11)
        desc_font.color.rgb = WHITE_COLOR

def create_j_curve_slide(prs):
    """Create J-curve slide"""
    slide = create_standard_slide(prs, "J-Curve & Capital Deployment Timeline")
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Investment Horizon and Return Realization"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    subtitle_font.bold = True
    
    # Timeline items
    timeline_items = [
        ("Anchor IPO Investments", "1-3 months post lock-in", "Shorter, shallower J-curve; quicker realization"),
        ("SME IPO Investments", "Medium-term", "Medium-term J-curve; exits via lock-in expiry or secondary market transactions"),
        ("Pre-IPO & Early-Stage Investments", "2-4 years", "Longer, potentially steeper J-curve; returns over 2-4 years, depending on company's IPO path or liquidity events")
    ]
    
    for i, (investment_type, duration, description) in enumerate(timeline_items):
        y_pos = Inches(2.8 + i * 1.3)
        
        # Timeline box
        timeline_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), y_pos, Inches(11.33), Inches(1))
        timeline_box.fill.solid()
        timeline_box.fill.fore_color.rgb = RGBColor(60, 60, 60)
        timeline_box.line.color.rgb = GOLD_COLOR
        timeline_box.line.width = Pt(2)
        
        # Investment type
        type_box = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.1), Inches(6), Inches(0.3))
        type_frame = type_box.text_frame
        type_frame.text = investment_type
        type_para = type_frame.paragraphs[0]
        type_font = type_para.font
        type_font.name = 'Segoe UI'
        type_font.size = Pt(16)
        type_font.color.rgb = GOLD_COLOR
        type_font.bold = True
        
        # Duration badge
        duration_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), y_pos + Inches(0.05), Inches(2.5), Inches(0.4))
        duration_box.fill.solid()
        duration_box.fill.fore_color.rgb = GOLD_COLOR
        duration_box.line.width = Pt(0)
        
        duration_text = slide.shapes.add_textbox(Inches(8.5), y_pos + Inches(0.05), Inches(2.5), Inches(0.4))
        duration_text_frame = duration_text.text_frame
        duration_text_frame.text = duration
        duration_text_para = duration_text_frame.paragraphs[0]
        duration_text_para.alignment = PP_ALIGN.CENTER
        duration_text_font = duration_text_para.font
        duration_text_font.name = 'Segoe UI'
        duration_text_font.size = Pt(12)
        duration_text_font.color.rgb = DARK_GRAY
        duration_text_font.bold = True
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.5), Inches(10), Inches(0.4))
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_para = desc_frame.paragraphs[0]
        desc_font = desc_para.font
        desc_font.name = 'Segoe UI'
        desc_font.size = Pt(12)
        desc_font.color.rgb = WHITE_COLOR

def create_management_team_slide(prs):
    """Create management team slide"""
    slide = create_standard_slide(prs, "Fund Management Team")
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Expertise Driving Credora Capital"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    subtitle_font.bold = True
    
    # Fund Manager section
    manager_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(11.33), Inches(1.8))
    manager_box.fill.solid()
    manager_box.fill.fore_color.rgb = RGBColor(50, 50, 50)
    manager_box.line.color.rgb = GOLD_COLOR
    manager_box.line.width = Pt(3)
    
    # Manager title
    manager_title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.7), Inches(10), Inches(0.3))
    manager_title_frame = manager_title_box.text_frame
    manager_title_frame.text = "Fund Manager & Lead Investment Strategist"
    manager_title_para = manager_title_frame.paragraphs[0]
    manager_title_font = manager_title_para.font
    manager_title_font.name = 'Segoe UI'
    manager_title_font.size = Pt(18)
    manager_title_font.color.rgb = GOLD_COLOR
    manager_title_font.bold = True
    
    # Manager name
    manager_name_box = slide.shapes.add_textbox(Inches(1.2), Inches(3), Inches(3), Inches(0.3))
    manager_name_frame = manager_name_box.text_frame
    manager_name_frame.text = "Shubham Gupta"
    manager_name_para = manager_name_frame.paragraphs[0]
    manager_name_font = manager_name_para.font
    manager_name_font.name = 'Segoe UI'
    manager_name_font.size = Pt(16)
    manager_name_font.color.rgb = GOLD_COLOR
    manager_name_font.bold = True
    
    # Manager description
    manager_desc_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.4), Inches(10.5), Inches(0.8))
    manager_desc_frame = manager_desc_box.text_frame
    manager_desc_frame.text = "Over 10 years of experience in primary/secondary markets, specializing in IPOs and SME equity. Leveraging deep market insight and analytical expertise for investment decisions, portfolio construction, and exit planning."
    manager_desc_para = manager_desc_frame.paragraphs[0]
    manager_desc_font = manager_desc_para.font
    manager_desc_font.name = 'Segoe UI'
    manager_desc_font.size = Pt(14)
    manager_desc_font.color.rgb = WHITE_COLOR
    
    # Advisors section
    advisors_title_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(10), Inches(0.3))
    advisors_title_frame = advisors_title_box.text_frame
    advisors_title_frame.text = "Advisors & Consultants - Distinguished panel with expertise in:"
    advisors_title_para = advisors_title_frame.paragraphs[0]
    advisors_title_font = advisors_title_para.font
    advisors_title_font.name = 'Segoe UI'
    advisors_title_font.size = Pt(16)
    advisors_title_font.color.rgb = GOLD_COLOR
    advisors_title_font.bold = True
    
    expertise_areas = [
        "Regulatory Compliance & SEBI Guidelines",
        "IPO Due Diligence & SME Market Analysis", 
        "Legal & Tax Structuring",
        "Financial Modelling & Valuation",
        "Investor Relations & Fundraising Strategies"
    ]
    
    add_bullet_points(slide, expertise_areas, Inches(5))

def create_team_expansion_slide(prs):
    """Create team expansion slide"""
    slide = create_standard_slide(prs, "Future Team Expansion & Partnerships")
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Building for Scalability & Deal Flow"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    subtitle_font.bold = True
    
    # Future team section
    future_team_points = [
        "Investment Analysts (financial due diligence, SME research)",
        "Compliance & Risk Officers",
        "Operations & Investor Relations Professionals"
    ]
    
    # Future team title
    future_title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(0.3))
    future_title_frame = future_title_box.text_frame
    future_title_frame.text = "Future Team (Next 12 Months):"
    future_title_para = future_title_frame.paragraphs[0]
    future_title_font = future_title_para.font
    future_title_font.name = 'Segoe UI'
    future_title_font.size = Pt(18)
    future_title_font.color.rgb = GOLD_COLOR
    future_title_font.bold = True
    
    add_bullet_points(slide, future_team_points, Inches(2.9))
    
    # Strategic collaborations
    collab_title_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(10), Inches(0.3))
    collab_title_frame = collab_title_box.text_frame
    collab_title_frame.text = "Strategic Collaborations - Actively exploring partnerships with:"
    collab_title_para = collab_title_frame.paragraphs[0]
    collab_title_font = collab_title_para.font
    collab_title_font.name = 'Segoe UI'
    collab_title_font.size = Pt(18)
    collab_title_font.color.rgb = GOLD_COLOR
    collab_title_font.bold = True
    
    collab_points = [
        "SME Industry Associations",
        "Investment Banks & IPO Brokers",
        "Legal & Accounting Firms (SME/IPO Specialized)",
        "Research & Analytics Providers"
    ]
    
    add_bullet_points(slide, collab_points, Inches(4.9))

def create_fees_terms_slide(prs):
    """Create fees and terms slide"""
    slide = create_standard_slide(prs, "Fees, Terms & Contact Information")
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Key Investment Details"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    subtitle_font.bold = True
    
    # Fee structure
    fees = [
        ("Management Fee", "3% per annum of committed capital (charged quarterly)"),
        ("Hurdle Rate", "12% per annum preferred return to LPs"),
        ("Performance Fee (Carry)", "To be decided/standard 20% over hurdle rate (optional)")
    ]
    
    for i, (fee_type, fee_desc) in enumerate(fees):
        y_pos = Inches(2.7 + i * 0.5)
        
        # Fee type
        fee_type_box = slide.shapes.add_textbox(Inches(1), y_pos, Inches(4), Inches(0.4))
        fee_type_frame = fee_type_box.text_frame
        fee_type_frame.text = fee_type + ":"
        fee_type_para = fee_type_frame.paragraphs[0]
        fee_type_font = fee_type_para.font
        fee_type_font.name = 'Segoe UI'
        fee_type_font.size = Pt(14)
        fee_type_font.color.rgb = GOLD_COLOR
        fee_type_font.bold = True
        
        # Fee description
        fee_desc_box = slide.shapes.add_textbox(Inches(5.5), y_pos, Inches(6.5), Inches(0.4))
        fee_desc_frame = fee_desc_box.text_frame
        fee_desc_frame.text = fee_desc
        fee_desc_para = fee_desc_frame.paragraphs[0]
        fee_desc_font = fee_desc_para.font
        fee_desc_font.name = 'Segoe UI'
        fee_desc_font.size = Pt(14)
        fee_desc_font.color.rgb = WHITE_COLOR
    
    # Terms
    terms = [
        ("Minimum Lock-in", "3 years from capital deployment"),
        ("Redemption", "Allowed post lock-in, subject to fund liquidity and exit events"),
        ("Reporting", "Quarterly NAV calculations, SEBI valuation guidelines, transparent investor reporting"),
        ("Audit & Compliance", "Annual audit by independent auditors")
    ]
    
    terms_title_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(10), Inches(0.3))
    terms_title_frame = terms_title_box.text_frame
    terms_title_frame.text = "Key Terms:"
    terms_title_para = terms_title_frame.paragraphs[0]
    terms_title_font = terms_title_para.font
    terms_title_font.name = 'Segoe UI'
    terms_title_font.size = Pt(16)
    terms_title_font.color.rgb = GOLD_COLOR
    terms_title_font.bold = True
    
    for i, (term_type, term_desc) in enumerate(terms):
        y_pos = Inches(4.9 + i * 0.4)
        
        term_box = slide.shapes.add_textbox(Inches(1), y_pos, Inches(11), Inches(0.3))
        term_frame = term_box.text_frame
        term_frame.text = f"• {term_type}: {term_desc}"
        term_para = term_frame.paragraphs[0]
        term_font = term_para.font
        term_font.name = 'Segoe UI'
        term_font.size = Pt(12)
        term_font.color.rgb = WHITE_COLOR

def create_investor_eligibility_slide(prs):
    """Create investor eligibility slide"""
    slide = create_standard_slide(prs, "Investor Eligibility & Commitment")
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Who Can Invest?"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = WHITE_COLOR
    subtitle_font.bold = True
    
    # Eligibility
    eligibility_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(0.8))
    eligibility_frame = eligibility_box.text_frame
    eligibility_frame.text = "Eligible Investors: Qualified Institutional Buyers (QIBs) and High Net Worth Individuals (HNIs) as per SEBI guidelines"
    eligibility_para = eligibility_frame.paragraphs[0]
    eligibility_font = eligibility_para.font
    eligibility_font.name = 'Segoe UI'
    eligibility_font.size = Pt(16)
    eligibility_font.color.rgb = WHITE_COLOR
    
    # Minimum investment highlight
    min_investment_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.5), Inches(5.33), Inches(1.2))
    min_investment_box.fill.solid()
    min_investment_box.fill.fore_color.rgb = GOLD_COLOR
    min_investment_box.line.width = Pt(0)
    
    # Amount
    amount_box = slide.shapes.add_textbox(Inches(4), Inches(3.6), Inches(5.33), Inches(0.6))
    amount_frame = amount_box.text_frame
    amount_frame.text = "₹1 Crore"
    amount_para = amount_frame.paragraphs[0]
    amount_para.alignment = PP_ALIGN.CENTER
    amount_font = amount_para.font
    amount_font.name = 'Segoe UI'
    amount_font.size = Pt(28)
    amount_font.color.rgb = DARK_GRAY
    amount_font.bold = True
    
    # Label
    min_label_box = slide.shapes.add_textbox(Inches(4), Inches(4.2), Inches(5.33), Inches(0.4))
    min_label_frame = min_label_box.text_frame
    min_label_frame.text = "Minimum Investment"
    min_label_para = min_label_frame.paragraphs[0]
    min_label_para.alignment = PP_ALIGN.CENTER
    min_label_font = min_label_para.font
    min_label_font.name = 'Segoe UI'
    min_label_font.size = Pt(14)
    min_label_font.color.rgb = DARK_GRAY
    min_label_font.bold = True
    
    # Communication section
    comm_title_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(10), Inches(0.3))
    comm_title_frame = comm_title_box.text_frame
    comm_title_frame.text = "Investor Communication - Detailed quarterly reports including:"
    comm_title_para = comm_title_frame.paragraphs[0]
    comm_title_font = comm_title_para.font
    comm_title_font.name = 'Segoe UI'
    comm_title_font.size = Pt(16)
    comm_title_font.color.rgb = GOLD_COLOR
    comm_title_font.bold = True
    
    comm_points = ["Portfolio updates", "NAV statements", "Market outlook"]
    add_bullet_points(slide, comm_points, Inches(5.6))

def create_thank_you_slide(prs):
    """Create thank you slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    set_slide_background(slide)
    
    # Add main logo (larger, centered)
    try:
        logo_path = '/app/credora_logo.jpeg'
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(5.5), Inches(1.8), width=Inches(2.5))
    except Exception as e:
        print(f"Could not add main logo: {e}")
    
    # Thank You title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You"
    
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_font = title_para.font
    title_font.name = 'Segoe UI'
    title_font.size = Pt(42)
    title_font.color.rgb = GOLD_COLOR
    title_font.bold = True
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "We look forward to building long-term value with you"
    
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Segoe UI'
    subtitle_font.size = Pt(20)
    subtitle_font.color.rgb = WHITE_COLOR

if __name__ == "__main__":
    print("🚀 Creating Credora Capital PowerPoint Presentation...")
    output_file = create_credora_presentation()
    print(f"✅ Presentation created successfully!")
    print(f"📁 File location: {output_file}")
    print(f"📊 Total slides: 15")
    print("🎯 Ready for investor meetings!")